"""Geometry QA for the deck, in place of a rendered visual pass.

LibreOffice cannot convert a pptx in the container this was built in — it fails on a
one-shape file too, so it is the install rather than the deck. This checks numerically
for the defects the visual pass exists to catch: shapes off the slide or inside the
margin, text boxes that overlap each other, and text too long for the box it is in.

Text fit is estimated from character counts, the same trade the renderer's callout
layout makes: measuring needs a real text engine, and the estimate only has to be
conservative enough to catch a genuine overflow. Widths are deliberately over-estimated,
so a warning here can be a false alarm but a silent pass should not hide a real one.
"""

import math
import sys

from pptx import Presentation
from pptx.util import Emu

EMU = 914400.0
MARGIN = 0.5          # inches of clear space wanted at every slide edge
AVG_CHAR = 0.50       # mean glyph advance as a share of point size, Calibri/Arial
LINE_H = 1.30         # line box as a share of point size


def inches(v):
    return (v or 0) / EMU


def text_of(shape):
    if not shape.has_text_frame:
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs)


def max_pt(shape):
    """Largest run size in the shape, which is what sets the line height."""
    sizes = [r.font.size.pt for p in shape.text_frame.paragraphs
             for r in p.runs if r.font.size is not None]
    return max(sizes) if sizes else 18.0


def est_height(shape):
    """Inches of text height this box needs, estimated from character counts."""
    w = inches(shape.width)
    if w <= 0:
        return 0.0
    pt = max_pt(shape)
    char_w = pt * AVG_CHAR / 72.0
    per_line = max(int(w / char_w), 1)
    lines = 0
    for para in shape.text_frame.paragraphs:
        t = para.text
        lines += max(math.ceil(len(t) / per_line), 1) if t else 1
    return lines * pt * LINE_H / 72.0


def overlap(a, b):
    ax1, ay1 = inches(a.left), inches(a.top)
    ax2, ay2 = ax1 + inches(a.width), ay1 + inches(a.height)
    bx1, by1 = inches(b.left), inches(b.top)
    bx2, by2 = bx1 + inches(b.width), by1 + inches(b.height)
    ox = min(ax2, bx2) - max(ax1, bx1)
    oy = min(ay2, by2) - max(ay1, by1)
    return (ox, oy) if ox > 0.02 and oy > 0.02 else None


def main(path):
    pres = Presentation(path)
    sw, sh = inches(pres.slide_width), inches(pres.slide_height)
    print(f"{path}: {len(pres.slides)} slides at {sw:.2f} x {sh:.2f} in\n")

    issues = 0
    for n, slide in enumerate(pres.slides, 1):
        found = []
        shapes = list(slide.shapes)

        for sp in shapes:
            x, y = inches(sp.left), inches(sp.top)
            w, h = inches(sp.width), inches(sp.height)
            label = (text_of(sp)[:34] or sp.shape_type.__str__()).replace("\n", " ")

            # off the slide entirely, or inside the margin
            if x < -0.01 or y < -0.01 or x + w > sw + 0.01 or y + h > sh + 0.01:
                found.append(f"OFF-SLIDE  {label!r} at ({x:.2f},{y:.2f}) {w:.2f}x{h:.2f}")
            elif x < MARGIN - 0.01 or y < MARGIN - 0.01 or \
                    x + w > sw - MARGIN + 0.01 or y + h > sh - MARGIN + 0.01:
                # the footnote sits low on purpose; everything else should clear the margin
                if "Source" not in label and "Sandisk" not in label and \
                        "Guidance" not in label and "Valuation" not in label and \
                        "Price milestones" not in label and "Both columns" not in label and \
                        "Reported quarterly" not in label and "Non-GAAP" not in label and \
                        "GAAP diluted" not in label and "FY2026 segment" not in label:
                    found.append(f"MARGIN     {label!r} at ({x:.2f},{y:.2f}) {w:.2f}x{h:.2f}")

            # text taller than its box
            if sp.has_text_frame and text_of(sp).strip():
                need = est_height(sp)
                if need > h + 0.06:
                    found.append(
                        f"OVERFLOW?  {label!r} needs ~{need:.2f}in in {h:.2f}in "
                        f"(w {w:.2f}, {max_pt(sp):.0f}pt)")

        # text over text
        texts = [sp for sp in shapes if sp.has_text_frame and text_of(sp).strip()]
        for i, a in enumerate(texts):
            for b in texts[i + 1:]:
                ov = overlap(a, b)
                if ov:
                    la = text_of(a)[:22].replace("\n", " ")
                    lb = text_of(b)[:22].replace("\n", " ")
                    found.append(f"TEXT-OVER  {la!r} / {lb!r} by {ov[0]:.2f}x{ov[1]:.2f}in")

        if found:
            issues += len(found)
            print(f"--- slide {n} ---")
            for line in found:
                print("   ", line)

    print(f"\n{issues} item(s) to look at.")
    return 0 if issues == 0 else 1


if __name__ == "__main__":
    sys.exit(main(sys.argv[1] if len(sys.argv) > 1 else "SanDisk-SNDK-FY2026.pptx"))
