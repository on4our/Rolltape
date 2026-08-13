"""Render an approximate picture of the deck, for the visual half of QA.

LibreOffice cannot open a pptx in this container (it fails on a one-shape file too), so
the usual convert-to-pdf route is unavailable. This reads the geometry back out of the
finished pptx with python-pptx and lays it out as HTML at the same coordinates, which is
enough to see the things a rendered page would show: balance, alignment, colour weight,
and text that does not fit.

What it is NOT: a renderer. Native charts are drawn as labelled placeholders showing
their series values, and font metrics are the browser's rather than PowerPoint's. Trust
it for layout, not for the last pixel.

    python preview_deck.py && npx playwright screenshot ...   (see __main__)
"""

import html
import sys

from pptx import Presentation
from pptx.util import Emu

EMU = 914400.0
PX = 96.0  # css pixels per inch


def inches(v):
    return (v or 0) / EMU


def solid_fill(shape):
    try:
        f = shape.fill
        # MSO_FILL_TYPE is an Enum rather than an IntEnum in current python-pptx, so
        # `f.type == 1` is False even for a solid fill. Compare the value.
        if f.type is not None and int(f.type) == 1:
            # RGBColor subclasses tuple, so %-formatting it unpacks the channels and
            # raises. str() gives the hex.
            return "#" + str(f.fore_color.rgb)
    except Exception:
        pass
    return None


def run_style(run):
    f = run.font
    bits = []
    if f.size is not None:
        bits.append(f"font-size:{f.size.pt:.1f}pt")
    if f.bold:
        bits.append("font-weight:700")
    try:
        if f.color is not None and f.color.rgb is not None:
            bits.append(f"color:#{f.color.rgb}")
    except Exception:
        pass
    if f.name:
        bits.append(f"font-family:{f.name},sans-serif")
    return ";".join(bits)


def chart_summary(shape):
    """Series values, so a placeholder still shows whether the numbers look right."""
    try:
        plot = shape.chart.plots[0]
        cats = [str(c) for c in plot.categories]
        vals = list(plot.series[0].values)
        return " · ".join(f"{c} {v:g}" for c, v in zip(cats, vals))
    except Exception:
        return "chart"


def render(path, out):
    pres = Presentation(path)
    sw, sh = inches(pres.slide_width), inches(pres.slide_height)
    parts = ["""<!doctype html><meta charset="utf-8"><style>
body{background:#333;margin:0;padding:18px;font-family:Calibri,Arial,sans-serif}
.slide{position:relative;margin:0 auto 22px;box-shadow:0 3px 14px #0008;overflow:hidden}
.n{position:absolute;left:-14px;top:0;color:#bbb;font:11px monospace}
.sh{position:absolute;box-sizing:border-box}
.tx{position:absolute;box-sizing:border-box;white-space:pre-wrap;line-height:1.28}
.ch{position:absolute;box-sizing:border-box;border:1px dashed #3B82F6aa;
    color:#9fb4d4;font:11px monospace;padding:6px;overflow:hidden}
</style>"""]

    for n, slide in enumerate(pres.slides, 1):
        bg = "#0B0E14"
        try:
            if int(slide.background.fill.type) == 1:
                bg = "#" + str(slide.background.fill.fore_color.rgb)
        except Exception:
            pass
        parts.append(f'<div class="slide" style="width:{sw * PX:.0f}px;'
                     f'height:{sh * PX:.0f}px;background:{bg}">'
                     f'<div class="n">{n}</div>')

        for shape in slide.shapes:
            x, y = inches(shape.left) * PX, inches(shape.top) * PX
            w, h = inches(shape.width) * PX, inches(shape.height) * PX
            box = f"left:{x:.0f}px;top:{y:.0f}px;width:{w:.0f}px;height:{h:.0f}px"

            if shape.has_chart:
                parts.append(f'<div class="ch" style="{box}">'
                             f'{html.escape(chart_summary(shape))}</div>')
                continue

            fill = solid_fill(shape)
            if fill and not shape.has_text_frame:
                radius = "50%" if "OVAL" in str(shape.shape_type) else "9px"
                parts.append(f'<div class="sh" style="{box};background:{fill};'
                             f'border-radius:{radius}"></div>')
                continue
            if fill:
                radius = "50%" if "OVAL" in str(shape.shape_type) else "9px"
                parts.append(f'<div class="sh" style="{box};background:{fill};'
                             f'border-radius:{radius}"></div>')

            if shape.has_text_frame and shape.text_frame.text.strip():
                inner = []
                for para in shape.text_frame.paragraphs:
                    if not para.runs:
                        inner.append("<br>")
                        continue
                    pa = ""
                    if para.alignment is not None and "CENTER" in str(para.alignment):
                        pa = "text-align:center;"
                    seg = "".join(
                        f'<span style="{run_style(r)}">{html.escape(r.text)}</span>'
                        for r in para.runs)
                    bullet = "• " if para.level == 0 and _has_bullet(para) else ""
                    inner.append(f'<div style="{pa}">{bullet}{seg}</div>')
                valign = "flex-start"
                try:
                    if shape.text_frame.vertical_anchor is not None and \
                            "MIDDLE" in str(shape.text_frame.vertical_anchor):
                        valign = "center"
                except Exception:
                    pass
                parts.append(
                    f'<div class="tx" style="{box};color:#E8ECF4;display:flex;'
                    f'flex-direction:column;justify-content:{valign}">'
                    f'{"".join(inner)}</div>')

        parts.append("</div>")

    with open(out, "w") as fh:
        fh.write("".join(parts))
    print(f"wrote {out} ({len(pres.slides)} slides at {sw * PX:.0f}x{sh * PX:.0f}px)")


def _has_bullet(para):
    xml = para._p.xml
    return "buChar" in xml or "buAutoNum" in xml


if __name__ == "__main__":
    render(sys.argv[1] if len(sys.argv) > 1 else "SanDisk-SNDK-FY2026.pptx",
           sys.argv[2] if len(sys.argv) > 2 else "preview.html")
